from pptx import Presentation
from pptx.util import Inches, Pt, Emu

prs = Presentation(r'c:\_arb\VSC_projects\LundaLoggern\Documentation\LundaLogger_2026-02-27.pptx')

for i, slide in enumerate(prs.slides):
    print(f'\n===== Slide {i+1} (layout: {slide.slide_layout.name}) =====')
    for shape in slide.shapes:
        if shape.has_text_frame:
            all_text = shape.text_frame.text.strip()
            if all_text:
                print(f'  [{shape.name}]:')
                for j, para in enumerate(shape.text_frame.paragraphs):
                    lvl = para.level
                    text = para.text
                    if text.strip():
                        font_info = ""
                        if para.runs:
                            r = para.runs[0]
                            font_info = f" (font={r.font.name}, size={r.font.size}, bold={r.font.bold})"
                        print(f'    L{lvl} [{j}]: "{text}"{font_info}')
