"""Check layout placeholder structure for the layouts we want to use."""
from pptx import Presentation
from lxml import etree

PPTX = r'c:\_arb\VSC_projects\LundaLoggern\Documentation\LundaLogger_2026-02-27 - 2nd-try.pptx'
prs = Presentation(PPTX)
NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'

# Check layouts: 3.1 (full-width bullets), 3.6 (two-column), 2.1 (section divider)
for target in ['3.1', '3.6', '2.1', '3.7']:
    for layout in prs.slide_layouts:
        if target in layout.name:
            print(f"\n{'='*60}")
            print(f"LAYOUT: '{layout.name}'")
            print(f"{'='*60}")
            for ph in layout.placeholders:
                print(f"  idx={ph.placeholder_format.idx:2d}  name='{ph.name}'  type={ph.placeholder_format.type}")
                if ph.has_text_frame:
                    xml = etree.tostring(ph._element, pretty_print=True).decode()
                    # Just show the lstStyle and body content
                    if 'Content' in ph.name or 'Title' in ph.name or 'Text Placeholder' in ph.name:
                        # Show the txBody
                        txBody = ph._element.find(f'{{{NS.replace("drawingml", "presentationml").replace("/main", "/2006/main")}}}txBody')
                        if txBody is None:
                            ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
                            txBody = ph._element.find(f'{{{ns_p}}}txBody')
                        print(f"    pos: left={ph.left}, top={ph.top}, w={ph.width}, h={ph.height}")
            break
