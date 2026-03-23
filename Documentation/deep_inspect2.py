"""Deep inspect template formatting: font, color, spacing, bullet XML for each placeholder."""
from pptx import Presentation
from pptx.util import Pt, Emu
from lxml import etree

PPTX = r'c:\_arb\VSC_projects\LundaLoggern\Documentation\LundaLogger_2026-02-27 - 2nd-try.pptx'
prs = Presentation(PPTX)

NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'

print(f"Slides: {len(prs.slides)}")
print(f"Size: {prs.slide_width/914400:.1f}\" x {prs.slide_height/914400:.1f}\"")

# Inspect each slide in detail
for i, slide in enumerate(prs.slides):
    print(f"\n{'='*70}")
    print(f"SLIDE {i+1} — layout: '{slide.slide_layout.name}'")
    print(f"{'='*70}")
    for shape in slide.shapes:
        if not shape.has_text_frame:
            print(f"  [{shape.name}] (no text frame) shape_type={shape.shape_type}")
            continue
        tf = shape.text_frame
        print(f"\n  [{shape.name}] word_wrap={tf.word_wrap} auto_size={tf.auto_size}")
        print(f"    pos: left={shape.left}, top={shape.top}, w={shape.width}, h={shape.height}")
        
        for j, para in enumerate(tf.paragraphs):
            text = para.text
            if not text.strip() and j > 5:
                continue
            
            print(f"    para[{j}] level={para.level} align={para.alignment}")
            
            # Check paragraph XML
            pPr = para._p.find(f'{{{NS}}}pPr')
            if pPr is not None:
                xml_str = etree.tostring(pPr, pretty_print=False).decode()
                if len(xml_str) < 500:
                    print(f"      pPr XML: {xml_str}")
                else:
                    print(f"      pPr XML (truncated): {xml_str[:500]}")
            
            for k, run in enumerate(para.runs):
                f = run.font
                try:
                    clr = f.color.rgb if f.color and f.color.type else None
                except:
                    clr = "theme/inherited"
                print(f"      run[{k}]: \"{run.text[:80]}\" "
                      f"font={f.name} sz={f.size} bold={f.bold} italic={f.italic} color={clr}")

# Inspect the key layout
print("\n\n" + "="*70)
print("LAYOUT '3.5 - Heading with bullet text and chart' — Content Placeholder XML")
print("="*70)
for layout in prs.slide_layouts:
    if '3.5' in layout.name:
        for ph in layout.placeholders:
            if 'Content' in ph.name:
                print(f"\n  [{ph.name}] idx={ph.placeholder_format.idx}")
                xml = etree.tostring(ph._element, pretty_print=True).decode()
                print(xml[:4000])

# Also look at 1.1 and 3.12 layouts
for layout_name_part in ['1.1', '3.12']:
    print(f"\n\n{'='*70}")
    print(f"LAYOUT containing '{layout_name_part}'")
    print("="*70)
    for layout in prs.slide_layouts:
        if layout_name_part in layout.name:
            for ph in layout.placeholders:
                if ph.has_text_frame and ('Title' in ph.name or 'Subtitle' in ph.name or 'Content' in ph.name or 'Text Placeholder' in ph.name):
                    xml = etree.tostring(ph._element, pretty_print=True).decode()
                    if len(xml) < 3000:
                        print(f"\n  [{ph.name}]:\n{xml}")
                    else:
                        print(f"\n  [{ph.name}]:\n{xml[:3000]}...")
