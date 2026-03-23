"""
Update LundaLogger PPTX — v2 fixed.
Uses the pristine copy as source. No slide deletion — avoids ZIP corruption.
Reuses the existing template slide, adds new ones, and reorders.
"""
from pptx import Presentation
from pptx.util import Pt

# Use the copy as pristine source, save to 2nd-try
SRC = r'c:\_arb\VSC_projects\LundaLoggern\Documentation\LundaLogger_2026-02-27 copy - Copy.pptx'
DST = r'c:\_arb\VSC_projects\LundaLoggern\Documentation\LundaLogger_2026-02-27 - 2nd-try.pptx'

prs = Presentation(SRC)

# ── Collect layouts ───────────────────────────────────────────────────
layout_by_name = {l.name: l for l in prs.slide_layouts}
L_CONTENT_W = layout_by_name['3.1 Content - Heading and bullet text on white bg']
L_CONTENT_S = layout_by_name['3.2 Content - Heding and bullet text on snow bg']
L_TWO_COL   = layout_by_name['3.6 - Heading with 2 column bullet text snow bg']
L_DIVIDER   = layout_by_name['2.1 - Section divider blue']

# ── Helpers ───────────────────────────────────────────────────────────

def ph(slide, idx):
    """Get placeholder by idx."""
    for s in slide.placeholders:
        if s.placeholder_format.idx == idx:
            return s
    return None

def set_title(slide, text):
    p = ph(slide, 0)
    if p and p.has_text_frame:
        p.text_frame.paragraphs[0].text = text

def set_sub(slide, text, idx=14):
    p = ph(slide, idx)
    if p and p.has_text_frame:
        p.text_frame.paragraphs[0].text = text

def bullets(slide, ph_idx, items):
    """Fill placeholder with items. Item = str | (str, level) | (str, level, bold)."""
    placeholder = ph(slide, ph_idx)
    if not placeholder:
        print(f"  WARNING: idx={ph_idx} not found on slide")
        return
    tf = placeholder.text_frame
    for i, item in enumerate(items):
        text = item if isinstance(item, str) else item[0]
        level = 0 if isinstance(item, str) else item[1]
        bold = False if isinstance(item, str) or len(item) < 3 else item[2]
        
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.level = level
        run = para.add_run()
        run.text = text
        if bold:
            run.font.bold = True

def move_to_end(zero_idx):
    """Move slide at 0-based index to end."""
    lst = prs.slides._sldIdLst
    el = lst[zero_idx]
    lst.remove(el)
    lst.append(el)

# ══════════════════════════════════════════════════════════════════════
# CURRENT STATE: Slide 0 = Title, Slide 1 = Template (3.5), Slide 2 = End
# PLAN: Repurpose slide 1 for "The Problem", add slides 2–12, move End last
# ══════════════════════════════════════════════════════════════════════

# ── Slide 1: Update title ────────────────────────────────────────────
print("1. Title slide")
s1 = prs.slides[0]
for shape in s1.shapes:
    nm = shape.name
    if nm == "Title 1" and shape.has_text_frame:
        shape.text_frame.paragraphs[0].text = "LundaLoggern"
    elif nm == "Subtitle 2" and shape.has_text_frame:
        shape.text_frame.paragraphs[0].text = "Smart Data Logging for Intensive Care Ventilators"
    elif "Text Placeholder 3" in nm and shape.has_text_frame:
        shape.text_frame.paragraphs[0].text = "\u00c5ke  |  2026-03-16"

# ── Slide 2: Repurpose the existing 3.5 slide for "The Problem" ──────
# Layout 3.5 has: Title idx=0, Sub-heading idx=14, Left content idx=19, Right content idx=20
print("2. The Problem (reusing template slide)")
s2 = prs.slides[1]
set_title(s2, "The Problem")
set_sub(s2, "Why ventilator data capture needs to be easier")
# Use the LEFT content area (idx mapped from slide: "Content Placeholder 9" = left)
# On this slide the placeholders may have different names but same idx from layout
left_ph_idx = None
right_ph_idx = None
for shape in s2.placeholders:
    idx = shape.placeholder_format.idx
    if idx == 19 or (shape.has_text_frame and shape.left < 5000000 and 'Content' in shape.name):
        left_ph_idx = idx
    elif idx == 20 or (shape.has_text_frame and shape.left >= 5000000 and 'Content' in shape.name):
        right_ph_idx = idx

# Try to find the actual content placeholder indices on this slide
print(f"   Left content idx: {left_ph_idx}, Right: {right_ph_idx}")
# Use whichever is available
content_idx = left_ph_idx or 19
bullets(s2, content_idx, [
    "Intensive care ventilators generate critical patient data \u2014 but capturing it consistently is hard",
    "Manual data recording is error-prone and time-consuming for clinical staff",
    "Existing solutions are often expensive, complex, or tied to proprietary hospital IT systems",
    ("How can we make ventilator data logging simple, portable, and independent?", 0, True),
])

# ── Slide 3: Introducing LundaLoggern ────────────────────────────────
print("3. Introducing LundaLoggern")
s = prs.slides.add_slide(L_CONTENT_S)
set_title(s, "Introducing LundaLoggern")
set_sub(s, "A plug-and-play data logger for SERVO ventilators")
bullets(s, 18, [
    "Small, self-contained device \u2014 no hospital network or IT integration needed",
    "Connects directly to the ventilator\u2019s CIE port",
    "Automatically records breath-by-breath data, settings, and waveforms to SD card",
    "Built-in colour display shows live status at a glance",
    "Built-in WiFi lets you download data from your phone or laptop",
])

# ── Slide 4: How It Works ────────────────────────────────────────────
print("4. How It Works")
s = prs.slides.add_slide(L_CONTENT_W)
set_title(s, "How It Works")
set_sub(s, "Four simple steps \u2014 no training required")
bullets(s, 18, [
    ("Step 1 \u2014 Plug in", 0, True),
    ("Connect LundaLoggern to the ventilator\u2019s CIE port", 1),
    ("Step 2 \u2014 Auto-detect", 0, True),
    ("The device identifies the ventilator and starts logging immediately", 1),
    ("Step 3 \u2014 Enable WiFi", 0, True),
    ("Press the button and scan the QR code on the display", 1),
    ("Step 4 \u2014 Download", 0, True),
    ("Open the web dashboard on your phone or laptop to review and download data files", 1),
])

# ── Slide 5: What Gets Logged ────────────────────────────────────────
print("5. What Gets Logged")
s = prs.slides.add_slide(L_CONTENT_S)
set_title(s, "What Gets Logged")
set_sub(s, "Comprehensive, configurable data capture")
bullets(s, 18, [
    ("Breath metrics", 0, True),
    ("Tidal volume, respiratory rate, minute ventilation, FiO\u2082, PEEP, peak pressure, and more", 1),
    ("Ventilator settings", 0, True),
    ("Mode, patient category, compliance compensation, I:E ratio", 1),
    ("Waveforms", 0, True),
    ("Flow, airway pressure, volume, CO\u2082, Edi \u2014 fully configurable", 1),
    "All data timestamped automatically using the ventilator\u2019s own clock",
    "Choose exactly which parameters to track via simple text files on the SD card",
])

# ── Slide 6: Key Features & Benefits (two-column) ────────────────────
print("6. Key Features & Benefits")
s = prs.slides.add_slide(L_TWO_COL)
set_title(s, "Key Features & Benefits")
set_sub(s, "Designed for simplicity, privacy, and portability")
bullets(s, 19, [
    ("Plug-and-play", 0, True),
    ("No setup, training, or IT support needed", 1),
    ("WiFi off by default", 0, True),
    ("Minimises RF interference in the ICU", 1),
    ("Portable & compact", 0, True),
    ("Fits in a pocket; move between ventilators", 1),
    ("SD card storage", 0, True),
    ("Standard, removable, unlimited capacity", 1),
])
bullets(s, 20, [
    ("QR code connectivity", 0, True),
    ("Connect to WiFi in seconds", 1),
    ("Configurable logging", 0, True),
    ("Adapt to research or clinical needs", 1),
    ("Live display", 0, True),
    ("Connection status, ventilator ID, timestamps", 1),
    ("No cloud dependency", 0, True),
    ("Data stays local \u2014 patient privacy ensured", 1),
])

# ── Slide 7: Target Markets & Use Cases (two-column) ─────────────────
print("7. Target Markets & Use Cases")
s = prs.slides.add_slide(L_TWO_COL)
set_title(s, "Target Markets & Use Cases")
set_sub(s, "")
bullets(s, 19, [
    ("Clinical / Neonatal ICU", 0, True),
    ("Continuous monitoring for quality assurance", 1),
    ("Research data collection during clinical studies", 1),
    ("Biomedical Engineering", 0, True),
    ("Performance verification and troubleshooting", 1),
    ("Long-term trend analysis during service calls", 1),
])
bullets(s, 20, [
    ("Medical Device R&D", 0, True),
    ("Bench testing and field validation of prototypes", 1),
    ("Rapid data capture without proprietary software", 1),
    ("Education & Training", 0, True),
    ("Teaching tool for respiratory therapy programs", 1),
    ("Live data demonstration during training sessions", 1),
])

# ── Slide 8: Competitive Advantages ──────────────────────────────────
print("8. Competitive Advantages")
s = prs.slides.add_slide(L_CONTENT_W)
set_title(s, "Competitive Advantages")
set_sub(s, "Why LundaLoggern stands out")
bullets(s, 18, [
    ("No recurring costs", 0, True),
    ("No licenses, subscriptions, or cloud fees", 1),
    ("Vendor-independent data", 0, True),
    ("Plain text files \u2014 open and portable", 1),
    ("No hospital IT involvement", 0, True),
    ("Standalone device with its own WiFi access point", 1),
    ("Privacy by design", 0, True),
    ("Data never leaves the device unless the user downloads it", 1),
    ("Low-cost hardware", 0, True),
    ("Built on the widely available ESP32-S3 platform", 1),
    ("Fast deployment", 0, True),
    ("Working in under a minute, zero configuration required", 1),
])

# ── Slide 9: Product Demo ────────────────────────────────────────────
print("9. Product Demo")
s = prs.slides.add_slide(L_CONTENT_S)
set_title(s, "Product Demo")
set_sub(s, "What you see on the device and on your phone")
bullets(s, 18, [
    ("TFT display shows:", 0, True),
    ("Ventilator type and serial number", 1),
    ("SD card and COM status indicators", 1),
    ("WiFi IP address and QR code", 1),
    ("Web dashboard on any device:", 0, True),
    ("File listing with timestamps", 1),
    ("One-tap download and delete", 1),
    ("Configuration viewer", 1),
])

# ── Slide 10: Technical Snapshot ──────────────────────────────────────
print("10. Technical Snapshot")
s = prs.slides.add_slide(L_CONTENT_W)
set_title(s, "Technical Snapshot")
set_sub(s, "Quick reference \u2014 key specs at a glance")
bullets(s, 18, [
    ("Hardware", 0, True),
    ("LilyGo T-Display S3 (ESP32-S3), SD card reader, RS232 interface", 1),
    ("Display", 0, True),
    ("320 \u00d7 170 px colour TFT", 1),
    ("Connectivity", 0, True),
    ("WiFi access point \u2014 no internet needed", 1),
    ("Protocol", 0, True),
    ("SERVO CIE interface (standard on all SERVO ventilators)", 1),
    ("Power", 0, True),
    ("USB-C \u2014 can be powered from the ventilator\u2019s USB port", 1),
    ("Software", 0, True),
    ("Open-source Arduino / PlatformIO firmware", 1),
])

# ── Slide 11: Roadmap ────────────────────────────────────────────────
print("11. Roadmap")
s = prs.slides.add_slide(L_CONTENT_S)
set_title(s, "Roadmap")
set_sub(s, "Future possibilities")
bullets(s, 18, [
    "WiFi auto-shutoff after 5 minutes of inactivity",
    "Bluetooth connectivity for direct mobile app integration",
    "Multi-ventilator support \u2014 expand beyond SERVO to other brands",
    "Cloud sync option (opt-in) for centralized research data collection",
    "Real-time trend graphs on the web dashboard",
    "Battery operation for transport and ambulance use",
    "Regulatory pathway \u2014 CE marking / FDA clearance for clinical deployment",
])

# ── Slide 12: Call to Action ─────────────────────────────────────────
print("12. Call to Action")
s = prs.slides.add_slide(L_CONTENT_W)
set_title(s, "Call to Action")
set_sub(s, "\u201cWe have a working prototype \u2014 let\u2019s bring it to market.\u201d")
bullets(s, 18, [
    ("What we need from marketing:", 0, True),
    ("Market sizing and customer segmentation feedback", 1),
    ("Input on branding, naming, and visual identity", 1),
    ("Channel strategy \u2014 direct sales, distributors, OEM partnerships?", 1),
    ("Pricing model validation", 1),
])

# ── Slide 13: Q&A (section divider) ──────────────────────────────────
print("13. Q&A")
s = prs.slides.add_slide(L_DIVIDER)
set_title(s, "Q & A")
set_sub(s, "Thank you!", idx=1)

# ── Reorder: move End slide (currently index 2) to the very end ──────
# After adding 12 new slides, deck is:
# [0]=Title, [1]=TheProblem(3.5), [2]=End, [3]=Introducing, ... [14]=Q&A
print("\nReordering: moving End slide to last position...")
move_to_end(2)  # end slide was at index 2

# ── Save ──────────────────────────────────────────────────────────────
prs.save(DST)
print(f"\nSaved: {DST}")
print(f"Total slides: {len(prs.slides)}")

# Verify
prs2 = Presentation(DST)
print(f"\nVerification ({len(prs2.slides)} slides):")
for i, slide in enumerate(prs2.slides):
    title = ""
    for shape in slide.shapes:
        if 'Title' in shape.name and shape.has_text_frame:
            t = shape.text_frame.text.strip()
            if t:
                title = t
                break
    print(f"  {i+1:2d}. [{slide.slide_layout.name}] {title}")
