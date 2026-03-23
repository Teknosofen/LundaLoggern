from pptx import Presentation
from pptx.util import Inches, Pt, Emu

prs = Presentation(r'c:\_arb\VSC_projects\LundaLoggern\Documentation\LundaLogger_2026-02-27.pptx')
print(f'Slide width: {prs.slide_width}, height: {prs.slide_height}')
print(f'Slide width inches: {prs.slide_width/914400}, height inches: {prs.slide_height/914400}')
print(f'Number of slides: {len(prs.slides)}')
print(f'Number of layouts: {len(prs.slide_layouts)}')
for i, layout in enumerate(prs.slide_layouts):
    print(f'  Layout {i}: {layout.name} (placeholders: {len(layout.placeholders)})')
    for ph in layout.placeholders:
        print(f'    ph idx={ph.placeholder_format.idx}, name={ph.name}, type={ph.placeholder_format.type}')
print()
for i, slide in enumerate(prs.slides):
    print(f'--- Slide {i+1} (layout: {slide.slide_layout.name}) ---')
    for shape in slide.shapes:
        print(f'  Shape: {shape.shape_type}, name={shape.name}, has_text={shape.has_text_frame}')
        if shape.has_text_frame:
            for j, para in enumerate(shape.text_frame.paragraphs):
                text = para.text[:120]
                if text.strip():
                    print(f'    para[{j}]: "{text}"')
        if shape.has_table:
            tbl = shape.table
            print(f'    TABLE: {len(tbl.rows)} rows x {len(tbl.columns)} cols')
            for r in range(len(tbl.rows)):
                row_text = []
                for c in range(len(tbl.columns)):
                    row_text.append(tbl.cell(r, c).text[:40])
                print(f'      row {r}: {row_text}')
