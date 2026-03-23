"""Deep inspect template formatting: font, color, spacing, bullet XML for each placeholder."""
from pptx import Presentation
from pptx.util import Pt, Emu
from lxml import etree

PPTX = r'c:\_arb\VSC_projects\LundaLoggern\Documentation\LundaLogger_2026-02-27 - 2nd-try.pptx'
prs = Presentation(PPTX)

print(f"Slides: {len(prs.slides)}")
print(f"Size: {prs.slide_width/914400:.1f}\" x {prs.slide_height/914400:.1f}\"")
print()

# Inspect each slide in detail
for i, slide in enumerate(prs.slides):
    print(f"\n{'='*70}")
    print(f"SLIDE {i+1} — layout: '{slide.slide_layout.name}'")
    print(f"{'='*70}")
    for shape in slide.shapes:
        if not shape.has_text_frame:
            print(f"  [{shape.name}] (no text frame) shape_type={shape.shape_type}")
            continue
        tf = shape.text_frame
        has_content = any(p.text.strip() for p in tf.paragraphs)
        print(f"\n  [{shape.name}] word_wrap={tf.word_wrap} auto_size={tf.auto_size}")
        print(f"    position: left={shape.left}, top={shape.top}, w={shape.width}, h={shape.height}")
        
        for j, para in enumerate(tf.paragraphs):
            text = para.text
            if not text.strip() and j > 2:
                continue  # skip many empty paras
            
            # Paragraph-level properties
            pf = para.paragraph_format
            print(f"    para[{j}] level={para.level} align={pf.alignment} "
                  f"spc_before={pf.space_before} spc_after={pf.space_after} "
                  f"line_spacing={pf.line_spacing}")
            
            # Check for bullet XML
            pPr = para._p.find('{http://schemas.openxmlformats.org/drawingml/2006/main}pPr')
            if pPr is not None:
                buNone = pPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}buNone')
                buChar = pPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}buChar')
                buAutoNum = pPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}buAutoNum')
                buFont = pPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}buFont')
                buSzPct = pPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}buSzPct')
                buClr = pPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}buClr')
                defRPr = pPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}defRPr')
                
                if buChar is not None:
                    print(f"      bullet: char='{buChar.get('char')}'")
                if buAutoNum is not None:
                    print(f"      bullet: autonum type='{buAutoNum.get('type')}'")
                if buNone is not None:
                    print(f"      bullet: NONE")
                if buFont is not None:
                    print(f"      buFont: {buFont.get('typeface')}")
                if buSzPct is not None:
                    print(f"      buSzPct: {buSzPct.get('val')}")
                if buClr is not None:
                    print(f"      buClr: {etree.tostring(buClr, pretty_print=True).decode()[:200]}")
                if defRPr is not None:
                    print(f"      defRPr: sz={defRPr.get('sz')} b={defRPr.get('b')}")
            
            # Run-level properties
            for k, run in enumerate(para.runs):
                f = run.font
                print(f"      run[{k}]: \"{run.text[:80]}\" "
                      f"font={f.name} sz={f.size} bold={f.bold} italic={f.italic} color={f.color.rgb if f.color and f.color.type else None}")

# Also inspect the layout placeholders in detail for Content Placeholder
print("\n\n" + "="*70)
print("LAYOUT DETAILS for '3.5 - Heading with bullet text and chart'")
print("="*70)
for layout in prs.slide_layouts:
    if '3.5' in layout.name:
        for ph in layout.placeholders:
            if ph.has_text_frame:
                print(f"\n  Layout PH: [{ph.name}] idx={ph.placeholder_format.idx}")
                tf = ph.text_frame
                for j, para in enumerate(tf.paragraphs):
                    pf = para.paragraph_format
                    print(f"    para[{j}] level={para.level} align={pf.alignment} "
                          f"spc_before={pf.space_before} spc_after={pf.space_after}")
                    pPr = para._p.find('{http://schemas.openxmlformats.org/drawingml/2006/main}pPr')
                    if pPr is not None:
                        print(f"      XML: {etree.tostring(pPr, pretty_print=False).decode()[:300]}")
                    for k, run in enumerate(para.runs):
                        f = run.font
                        print(f"      run[{k}]: \"{run.text[:50]}\" font={f.name} sz={f.size} bold={f.bold} color={f.color.rgb if f.color and f.color.type else None}")

# Also check the "Content Placeholder 2" XML from the copy template  
print("\n\n" + "="*70)
print("RAW XML of Content Placeholder from layout")
print("="*70)
for layout in prs.slide_layouts:
    if '3.5' in layout.name:
        for ph in layout.placeholders:
            if 'Content' in ph.name:
                print(f"\n  [{ph.name}] idx={ph.placeholder_format.idx}")
                xml = etree.tostring(ph._element, pretty_print=True).decode()
                print(xml[:3000])
