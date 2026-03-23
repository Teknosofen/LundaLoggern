"""
Update LundaLogger_2026-02-27 - 2nd-try.pptx with content from PresentationOutline.md
Preserves all template formatting — no explicit font/size/color overrides.
"""
from pptx import Presentation
from pptx.util import Pt, Emu
from copy import deepcopy
from lxml import etree

PPTX_PATH = r'c:\_arb\VSC_projects\LundaLoggern\Documentation\LundaLogger_2026-02-27 - 2nd-try.pptx'

prs = Presentation(PPTX_PATH)

# ── Layout references ─────────────────────────────────────────────────
layouts = {name: layout for layout in prs.slide_layouts for name in [layout.name]}
layout_map = {}
for layout in prs.slide_layouts:
    layout_map[layout.name] = layout

L_CONTENT     = layout_map['3.1 Content - Heading and bullet text on white bg']
L_CONTENT_SNO = layout_map['3.2 Content - Heding and bullet text on snow bg']
L_TWO_COL     = layout_map['3.6 - Heading with 2 column bullet text snow bg']
L_DIVIDER     = layout_map['2.1 - Section divider blue']
L_CONTENT_PIC = layout_map['3.5 - Heading with bullet text and chart']
L_END         = layout_map['4.2 - End logo on graphic']

# ── Helpers ───────────────────────────────────────────────────────────

def find_ph(slide, idx):
    """Find a placeholder by idx on a slide."""
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == idx:
            return shape
    return None


def set_title(slide, text):
    """Set the title placeholder (idx=0)."""
    ph = find_ph(slide, 0)
    if ph and ph.has_text_frame:
        ph.text_frame.paragraphs[0].text = text


def set_subtitle(slide, text, idx=14):
    """Set the sub-heading placeholder (idx=14 on 3.x layouts, idx=1 on 2.x)."""
    ph = find_ph(slide, idx)
    if ph and ph.has_text_frame:
        ph.text_frame.paragraphs[0].text = text


def fill_bullets(slide, ph_idx, items):
    """
    Fill a content placeholder with bullet items.
    
    Each item is either:
      - A string (level 0, normal)
      - A tuple: (text, level)           — normal text at given level
      - A tuple: (text, level, True)      — bold run at given level
    
    NO explicit font properties are set — the template theme handles all styling.
    """
    ph = find_ph(slide, ph_idx)
    if ph is None:
        print(f"  WARNING: placeholder idx={ph_idx} not found")
        return
    
    tf = ph.text_frame
    
    # The first paragraph already exists from the layout; reuse it
    for i, item in enumerate(items):
        if isinstance(item, str):
            text, level, bold = item, 0, False
        elif len(item) == 2:
            text, level = item
            bold = False
        else:
            text, level, bold = item
        
        if i == 0:
            para = tf.paragraphs[0]
        else:
            para = tf.add_paragraph()
        
        para.level = level
        run = para.add_run()
        run.text = text
        if bold:
            run.font.bold = True


def add_slide(layout):
    """Add a new slide at the end (before we rearrange)."""
    return prs.slides.add_slide(layout)


def move_slide_to_end(slide_index):
    """Move a slide (by 0-based index) to the end of the deck."""
    sldIdLst = prs.slides._sldIdLst
    el = sldIdLst[slide_index]
    sldIdLst.remove(el)
    sldIdLst.append(el)


def delete_slide(slide_index):
    """Delete slide at 0-based index."""
    rId = prs.slides._sldIdLst[slide_index].get(
        '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id'
    )
    prs.part.drop_rel(rId)
    slide_elem = prs.slides._sldIdLst[slide_index]
    prs.slides._sldIdLst.remove(slide_elem)


# ══════════════════════════════════════════════════════════════════════
# BUILD THE PRESENTATION
# ══════════════════════════════════════════════════════════════════════

# ── Slide 1: Title (already exists — update text) ────────────────────
print("Slide 1: Title")
slide1 = prs.slides[0]
# Title
for shape in slide1.shapes:
    if shape.name == "Title 1" and shape.has_text_frame:
        shape.text_frame.paragraphs[0].text = "LundaLoggern"
    elif shape.name == "Subtitle 2" and shape.has_text_frame:
        shape.text_frame.paragraphs[0].text = "Smart Data Logging for Intensive Care Ventilators"
    elif shape.name == "Text Placeholder 3" and shape.has_text_frame:
        shape.text_frame.paragraphs[0].text = "\u00c5ke  |  2026-03-16"

# ── Remove the empty template slide (index 1) ────────────────────────
print("Removing empty template slide...")
delete_slide(1)
# Now we have: [Title, End]. New slides will be inserted between.

# ── Slide 2: The Problem (3.1 white bg) ──────────────────────────────
print("Slide 2: The Problem")
s = add_slide(L_CONTENT)
set_title(s, "The Problem")
set_subtitle(s, "Why ventilator data capture needs to be easier")
fill_bullets(s, 18, [
    "Intensive care ventilators generate critical patient data \u2014 but capturing it consistently is hard",
    "Manual data recording is error-prone and time-consuming for clinical staff",
    "Existing solutions are often expensive, complex, or tied to proprietary hospital IT systems",
    ("How can we make ventilator data logging simple, portable, and independent?", 0, True),
])

# ── Slide 3: Introducing LundaLoggern (3.1 snow bg for variety) ──────
print("Slide 3: Introducing LundaLoggern")
s = add_slide(L_CONTENT_SNO)
set_title(s, "Introducing LundaLoggern")
set_subtitle(s, "A plug-and-play data logger for SERVO ventilators")
fill_bullets(s, 18, [
    "Small, self-contained device \u2014 no hospital network or IT integration needed",
    "Connects directly to the ventilator\u2019s CIE port",
    "Automatically records breath-by-breath data, settings, and waveforms to SD card",
    "Built-in colour display shows live status at a glance",
    "Built-in WiFi lets you download data from your phone or laptop",
])

# ── Slide 4: How It Works (3.1 white bg) ─────────────────────────────
print("Slide 4: How It Works")
s = add_slide(L_CONTENT)
set_title(s, "How It Works")
set_subtitle(s, "Four simple steps \u2014 no training required")
fill_bullets(s, 18, [
    ("Step 1 \u2014 Plug in", 0, True),
    ("Connect LundaLoggern to the ventilator\u2019s CIE port", 1),
    ("Step 2 \u2014 Auto-detect", 0, True),
    ("The device identifies the ventilator and starts logging immediately", 1),
    ("Step 3 \u2014 Enable WiFi", 0, True),
    ("Press the button and scan the QR code on the display", 1),
    ("Step 4 \u2014 Download", 0, True),
    ("Open the web dashboard on your phone or laptop to review and download data files", 1),
])

# ── Slide 5: What Gets Logged (3.1 snow bg) ──────────────────────────
print("Slide 5: What Gets Logged")
s = add_slide(L_CONTENT_SNO)
set_title(s, "What Gets Logged")
set_subtitle(s, "Comprehensive, configurable data capture")
fill_bullets(s, 18, [
    ("Breath metrics", 0, True),
    ("Tidal volume, respiratory rate, minute ventilation, FiO\u2082, PEEP, peak pressure, and more", 1),
    ("Ventilator settings", 0, True),
    ("Mode, patient category, compliance compensation, I:E ratio", 1),
    ("Waveforms", 0, True),
    ("Flow, airway pressure, volume, CO\u2082, Edi \u2014 fully configurable", 1),
    "All data timestamped automatically using the ventilator\u2019s own clock",
    "Choose exactly which parameters to track via simple text files on the SD card",
])

# ── Slide 6: Key Features & Benefits (3.6 two-column) ────────────────
print("Slide 6: Key Features & Benefits")
s = add_slide(L_TWO_COL)
set_title(s, "Key Features & Benefits")
set_subtitle(s, "Designed for simplicity, privacy, and portability")
# Left column (idx=19)
fill_bullets(s, 19, [
    ("Plug-and-play", 0, True),
    ("No setup, training, or IT support needed", 1),
    ("WiFi off by default", 0, True),
    ("Minimises RF interference in the ICU", 1),
    ("Portable & compact", 0, True),
    ("Fits in a pocket; move between ventilators", 1),
    ("SD card storage", 0, True),
    ("Standard, removable, unlimited capacity", 1),
])
# Right column (idx=20)
fill_bullets(s, 20, [
    ("QR code connectivity", 0, True),
    ("Connect to WiFi in seconds", 1),
    ("Configurable logging", 0, True),
    ("Adapt to research or clinical needs", 1),
    ("Live display", 0, True),
    ("Connection status, ventilator ID, timestamps", 1),
    ("No cloud dependency", 0, True),
    ("Data stays local \u2014 patient privacy ensured", 1),
])

# ── Slide 7: Target Markets (3.6 two-column) ─────────────────────────
print("Slide 7: Target Markets & Use Cases")
s = add_slide(L_TWO_COL)
set_title(s, "Target Markets & Use Cases")
set_subtitle(s, "")
# Left column
fill_bullets(s, 19, [
    ("Clinical / Neonatal ICU", 0, True),
    ("Continuous monitoring for quality assurance", 1),
    ("Research data collection during clinical studies", 1),
    ("Biomedical Engineering", 0, True),
    ("Performance verification and troubleshooting", 1),
    ("Long-term trend analysis during service calls", 1),
])
# Right column
fill_bullets(s, 20, [
    ("Medical Device R&D", 0, True),
    ("Bench testing and field validation of prototypes", 1),
    ("Rapid data capture without proprietary software", 1),
    ("Education & Training", 0, True),
    ("Teaching tool for respiratory therapy programs", 1),
    ("Live data demonstration during training sessions", 1),
])

# ── Slide 8: Competitive Advantages (3.1 white bg) ───────────────────
print("Slide 8: Competitive Advantages")
s = add_slide(L_CONTENT)
set_title(s, "Competitive Advantages")
set_subtitle(s, "Why LundaLoggern stands out")
fill_bullets(s, 18, [
    ("No recurring costs", 0, True),
    ("No licenses, subscriptions, or cloud fees", 1),
    ("Vendor-independent data", 0, True),
    ("Plain text files \u2014 open and portable", 1),
    ("No hospital IT involvement", 0, True),
    ("Standalone device with its own WiFi access point", 1),
    ("Privacy by design", 0, True),
    ("Data never leaves the device unless the user downloads it", 1),
    ("Low-cost hardware", 0, True),
    ("Built on the widely available ESP32-S3 platform", 1),
    ("Fast deployment", 0, True),
    ("Working in under a minute, zero configuration required", 1),
])

# ── Slide 9: Product Demo (3.1 snow bg) ──────────────────────────────
print("Slide 9: Product Demo / Live Screenshots")
s = add_slide(L_CONTENT_SNO)
set_title(s, "Product Demo")
set_subtitle(s, "What you see on the device and on your phone")
fill_bullets(s, 18, [
    ("TFT display shows:", 0, True),
    ("Ventilator type and serial number", 1),
    ("SD card and COM status indicators", 1),
    ("WiFi IP address and QR code", 1),
    ("Web dashboard on any device:", 0, True),
    ("File listing with timestamps", 1),
    ("One-tap download and delete", 1),
    ("Configuration viewer", 1),
])

# ── Slide 10: Technical Snapshot (3.1 white bg) ──────────────────────
print("Slide 10: Technical Snapshot")
s = add_slide(L_CONTENT)
set_title(s, "Technical Snapshot")
set_subtitle(s, "Quick reference \u2014 key specs at a glance")
fill_bullets(s, 18, [
    ("Hardware", 0, True),
    ("LilyGo T-Display S3 (ESP32-S3), SD card reader, RS232 interface", 1),
    ("Display", 0, True),
    ("320 \u00d7 170 px colour TFT", 1),
    ("Connectivity", 0, True),
    ("WiFi access point \u2014 no internet needed", 1),
    ("Protocol", 0, True),
    ("SERVO CIE interface (standard on all SERVO ventilators)", 1),
    ("Power", 0, True),
    ("USB-C \u2014 can be powered from the ventilator\u2019s USB port", 1),
    ("Software", 0, True),
    ("Open-source Arduino / PlatformIO firmware", 1),
])

# ── Slide 11: Roadmap (3.1 snow bg) ──────────────────────────────────
print("Slide 11: Roadmap / Future Possibilities")
s = add_slide(L_CONTENT_SNO)
set_title(s, "Roadmap")
set_subtitle(s, "Future possibilities")
fill_bullets(s, 18, [
    "WiFi auto-shutoff after 5 minutes of inactivity",
    "Bluetooth connectivity for direct mobile app integration",
    "Multi-ventilator support \u2014 expand beyond SERVO to other brands",
    "Cloud sync option (opt-in) for centralized research data collection",
    "Real-time trend graphs on the web dashboard",
    "Battery operation for transport and ambulance use",
    "Regulatory pathway \u2014 CE marking / FDA clearance for clinical deployment",
])

# ── Slide 12: Call to Action (3.1 white bg) ──────────────────────────
print("Slide 12: Call to Action")
s = add_slide(L_CONTENT)
set_title(s, "Call to Action")
set_subtitle(s, "\u201cWe have a working prototype \u2014 let\u2019s bring it to market.\u201d")
fill_bullets(s, 18, [
    ("What we need from marketing:", 0, True),
    ("Market sizing and customer segmentation feedback", 1),
    ("Input on branding, naming, and visual identity", 1),
    ("Channel strategy \u2014 direct sales, distributors, OEM partnerships?", 1),
    ("Pricing model validation", 1),
])

# ── Slide 13: Q&A (section divider blue) ─────────────────────────────
print("Slide 13: Q&A")
s = add_slide(L_DIVIDER)
set_title(s, "Q & A")
set_subtitle(s, "Thank you!", idx=1)

# ── Move the End slide to the very end ────────────────────────────────
# The End slide is currently at index 1 (right after the title).
# All new slides were appended after it. We need to move it to the end.
print("Moving end slide to last position...")
move_slide_to_end(1)

# ── Save ──────────────────────────────────────────────────────────────
prs.save(PPTX_PATH)
print(f"\nDone! Saved to: {PPTX_PATH}")
print(f"Total slides: {len(prs.slides)}")

# Quick verification
for i, slide in enumerate(prs.slides):
    title_text = ""
    for shape in slide.shapes:
        if shape.has_text_frame:
            t = shape.text_frame.text.strip()
            if t and ('Title' in shape.name):
                title_text = t
                break
    print(f"  Slide {i+1}: {slide.slide_layout.name} — \"{title_text}\"")
