"""
Update LundaLogger PPTX — v2 final.
Source: pristine 'copy - Copy' file. Output: '2nd-try' file.
No slide deletion. Moves the blank template slide out of the way.
"""
from pptx import Presentation

SRC = r'c:\_arb\VSC_projects\LundaLoggern\Documentation\LundaLogger_2026-02-27 -3rd-try.pptx'
DST = r'c:\_arb\VSC_projects\LundaLoggern\Documentation\LundaLogger_2026-02-27 - 2nd-try.pptx'

prs = Presentation(SRC)

# ── Layouts ───────────────────────────────────────────────────────────
L = {l.name: l for l in prs.slide_layouts}
L_W = L['3.1 Content - Heading and bullet text on white bg']
L_S = L['3.2 Content - Heding and bullet text on snow bg']
L_2C = L['3.6 - Heading with 2 column bullet text snow bg']
L_DIV = L['2.1 - Section divider blue']

# ── Helpers ───────────────────────────────────────────────────────────

def ph(slide, idx):
    for s in slide.placeholders:
        if s.placeholder_format.idx == idx:
            return s
    return None

def title(slide, text):
    p = ph(slide, 0)
    if p and p.has_text_frame:
        p.text_frame.paragraphs[0].text = text

def sub(slide, text, idx=14):
    p = ph(slide, idx)
    if p and p.has_text_frame:
        p.text_frame.paragraphs[0].text = text

def bul(slide, ph_idx, items):
    """Items: str | (str,level) | (str,level,bold)."""
    p = ph(slide, ph_idx)
    if not p or not p.has_text_frame:
        print(f"  WARN: idx={ph_idx} missing or no text frame")
        return
    tf = p.text_frame
    for i, item in enumerate(items):
        text = item if isinstance(item, str) else item[0]
        level = 0 if isinstance(item, str) else item[1]
        bold = False if (isinstance(item, str) or len(item) < 3) else item[2]
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.level = level
        run = para.add_run()
        run.text = text
        if bold:
            run.font.bold = True

def move_to_end(zero_idx):
    lst = prs.slides._sldIdLst
    el = lst[zero_idx]
    lst.remove(el)
    lst.append(el)

# ══════════════════════════════════════════════════════════════════════
# STATE: [0]=Title(1.1), [1]=Template(3.12), [2]=Template(3.5), [3]=End(4.2)
# PLAN: Update [0], add 12 new slides, move [1],[2],[3] to end
# ══════════════════════════════════════════════════════════════════════

# ── 1. Title slide (update in place) ─────────────────────────────────
print(" 1. Title")
s0 = prs.slides[0]
for shape in s0.shapes:
    n = shape.name
    if n == "Title 1" and shape.has_text_frame:
        shape.text_frame.paragraphs[0].text = "LundaLoggern"
    elif n == "Subtitle 2" and shape.has_text_frame:
        shape.text_frame.paragraphs[0].text = "Smart Data Logging for Intensive Care Ventilators"
    elif "Text Placeholder 3" in n and shape.has_text_frame:
        shape.text_frame.paragraphs[0].text = "\u00c5ke  |  2026-03-16"

# ── 2. The Problem ───────────────────────────────────────────────────
print(" 2. The Problem")
s = prs.slides.add_slide(L_W)
title(s, "The Problem")
sub(s, "Why ventilator data capture needs to be easier")
bul(s, 18, [
    "Intensive care ventilators generate critical patient data \u2014 but capturing it consistently is hard",
    "Manual data recording is error-prone and time-consuming for clinical staff",
    "Existing solutions are often expensive, complex, or tied to proprietary hospital IT systems",
    ("How can we make ventilator data logging simple, portable, and independent?", 0, True),
])

# ── 3. Introducing LundaLoggern ──────────────────────────────────────
print(" 3. Introducing LundaLoggern")
s = prs.slides.add_slide(L_S)
title(s, "Introducing LundaLoggern")
sub(s, "A plug-and-play data logger for SERVO ventilators")
bul(s, 18, [
    "Small, self-contained device \u2014 no hospital network or IT integration needed",
    "Connects directly to the ventilator\u2019s CIE port",
    "Automatically records breath-by-breath data, settings, and waveforms to SD card",
    "Built-in colour display shows live status at a glance",
    "Built-in WiFi lets you download data from your phone or laptop",
])

# ── 4. How It Works ──────────────────────────────────────────────────
print(" 4. How It Works")
s = prs.slides.add_slide(L_W)
title(s, "How It Works")
sub(s, "Four simple steps \u2014 no training required")
bul(s, 18, [
    ("Step 1 \u2014 Plug in", 0, True),
    ("Connect LundaLoggern to the ventilator\u2019s CIE port", 1),
    ("Step 2 \u2014 Auto-detect", 0, True),
    ("The device identifies the ventilator and starts logging immediately", 1),
    ("Step 3 \u2014 Enable WiFi", 0, True),
    ("Press the button and scan the QR code on the display", 1),
    ("Step 4 \u2014 Download", 0, True),
    ("Open the web dashboard on your phone or laptop to review and download data files", 1),
])

# ── 5. What Gets Logged ──────────────────────────────────────────────
print(" 5. What Gets Logged")
s = prs.slides.add_slide(L_S)
title(s, "What Gets Logged")
sub(s, "Comprehensive, configurable data capture")
bul(s, 18, [
    ("Breath metrics", 0, True),
    ("Tidal volume, respiratory rate, minute ventilation, FiO\u2082, PEEP, peak pressure, and more", 1),
    ("Ventilator settings", 0, True),
    ("Mode, patient category, compliance compensation, I:E ratio", 1),
    ("Waveforms", 0, True),
    ("Flow, airway pressure, volume, CO\u2082, Edi \u2014 fully configurable", 1),
    "All data timestamped automatically using the ventilator\u2019s own clock",
    "Choose which parameters to track via simple text files on the SD card",
])

# ── 6. Key Features & Benefits (two-column) ──────────────────────────
print(" 6. Key Features & Benefits")
s = prs.slides.add_slide(L_2C)
title(s, "Key Features & Benefits")
sub(s, "Designed for simplicity, privacy, and portability")
bul(s, 19, [
    ("Plug-and-play", 0, True),
    ("No setup, training, or IT support needed", 1),
    ("WiFi off by default", 0, True),
    ("Minimises RF interference in the ICU", 1),
    ("Portable & compact", 0, True),
    ("Fits in a pocket; move between ventilators", 1),
    ("SD card storage", 0, True),
    ("Standard, removable, unlimited capacity", 1),
])
bul(s, 20, [
    ("QR code connectivity", 0, True),
    ("Connect to WiFi in seconds", 1),
    ("Configurable logging", 0, True),
    ("Adapt to research or clinical needs", 1),
    ("Live display", 0, True),
    ("Connection status, ventilator ID, timestamps", 1),
    ("No cloud dependency", 0, True),
    ("Data stays local \u2014 patient privacy ensured", 1),
])

# ── 7. Target Markets (two-column) ───────────────────────────────────
print(" 7. Target Markets & Use Cases")
s = prs.slides.add_slide(L_2C)
title(s, "Target Markets & Use Cases")
sub(s, "")
bul(s, 19, [
    ("Clinical / Neonatal ICU", 0, True),
    ("Continuous monitoring for quality assurance", 1),
    ("Research data collection during clinical studies", 1),
    ("Biomedical Engineering", 0, True),
    ("Performance verification and troubleshooting", 1),
    ("Long-term trend analysis during service calls", 1),
])
bul(s, 20, [
    ("Medical Device R&D", 0, True),
    ("Bench testing and field validation of prototypes", 1),
    ("Rapid data capture without proprietary software", 1),
    ("Education & Training", 0, True),
    ("Teaching tool for respiratory therapy programs", 1),
    ("Live data demonstration during training sessions", 1),
])

# ── 8. Competitive Advantages ────────────────────────────────────────
print(" 8. Competitive Advantages")
s = prs.slides.add_slide(L_W)
title(s, "Competitive Advantages")
sub(s, "Why LundaLoggern stands out")
bul(s, 18, [
    ("No recurring costs", 0, True),
    ("No licenses, subscriptions, or cloud fees", 1),
    ("Vendor-independent data", 0, True),
    ("Plain text files \u2014 open and portable", 1),
    ("No hospital IT involvement", 0, True),
    ("Standalone device with its own WiFi access point", 1),
    ("Privacy by design", 0, True),
    ("Data never leaves the device unless the user downloads it", 1),
    ("Low-cost hardware", 0, True),
    ("Built on the widely available ESP32-S3 platform", 1),
    ("Fast deployment", 0, True),
    ("Working in under a minute, zero configuration required", 1),
])

# ── 9. Product Demo ──────────────────────────────────────────────────
print(" 9. Product Demo")
s = prs.slides.add_slide(L_S)
title(s, "Product Demo")
sub(s, "What you see on the device and on your phone")
bul(s, 18, [
    ("TFT display shows:", 0, True),
    ("Ventilator type and serial number", 1),
    ("SD card and COM status indicators", 1),
    ("WiFi IP address and QR code", 1),
    ("Web dashboard on any device:", 0, True),
    ("File listing with timestamps", 1),
    ("One-tap download and delete", 1),
    ("Configuration viewer", 1),
])

# ── 10. Technical Snapshot ────────────────────────────────────────────
print("10. Technical Snapshot")
s = prs.slides.add_slide(L_W)
title(s, "Technical Snapshot")
sub(s, "Quick reference \u2014 key specs at a glance")
bul(s, 18, [
    ("Hardware", 0, True),
    ("LilyGo T-Display S3 (ESP32-S3), SD card reader, RS232 interface", 1),
    ("Display", 0, True),
    ("320 \u00d7 170 px colour TFT", 1),
    ("Connectivity", 0, True),
    ("WiFi access point \u2014 no internet needed", 1),
    ("Protocol", 0, True),
    ("SERVO CIE interface (standard on all SERVO ventilators)", 1),
    ("Power", 0, True),
    ("USB-C \u2014 can be powered from the ventilator\u2019s USB port", 1),
    ("Software", 0, True),
    ("Open-source Arduino / PlatformIO firmware", 1),
])

# ── 11. Roadmap ──────────────────────────────────────────────────────
print("11. Roadmap")
s = prs.slides.add_slide(L_S)
title(s, "Roadmap")
sub(s, "Future possibilities")
bul(s, 18, [
    "WiFi auto-shutoff after 5 minutes of inactivity",
    "Bluetooth connectivity for direct mobile app integration",
    "Multi-ventilator support \u2014 expand beyond SERVO to other brands",
    "Cloud sync option (opt-in) for centralized research data collection",
    "Real-time trend graphs on the web dashboard",
    "Battery operation for transport and ambulance use",
    "Regulatory pathway \u2014 CE marking / FDA clearance for clinical deployment",
])

# ── 12. Call to Action ───────────────────────────────────────────────
print("12. Call to Action")
s = prs.slides.add_slide(L_W)
title(s, "Call to Action")
sub(s, "\u201cWe have a working prototype \u2014 let\u2019s bring it to market.\u201d")
bul(s, 18, [
    ("What we need from marketing:", 0, True),
    ("Market sizing and customer segmentation feedback", 1),
    ("Input on branding, naming, and visual identity", 1),
    ("Channel strategy \u2014 direct sales, distributors, OEM partnerships?", 1),
    ("Pricing model validation", 1),
])

# ── 13. Q&A ──────────────────────────────────────────────────────────
print("13. Q&A")
s = prs.slides.add_slide(L_DIV)
title(s, "Q & A")
sub(s, "Thank you!", idx=1)

# ══════════════════════════════════════════════════════════════════════
# REORDER: Current order is:
#   [0]=Title, [1]=Template(3.12), [2]=Template(3.5), [3]=End(4.2), [4..15]=content
# Target order:
#   [0]=Title, [4..15]=content, [3]=End, [1],[2]=leftover templates
# ══════════════════════════════════════════════════════════════════════
print("\nReordering slides...")

# Move template slides (indices 1 and 2) and end slide (index 3) to the end
# Move index 1 to end three times (each time the next template shifts down)
move_to_end(1)  # Template(3.12) -> end
# Now: [0]=Title, [1]=Template(3.5), [2]=End, [3..14]=content, [14]=Template(3.12)
move_to_end(1)  # Template(3.5) -> end
# Now: [0]=Title, [1]=End, [2..13]=content, [14]=Template(3.12), [15]=Template(3.5)
move_to_end(1)  # End -> end
# Now: [0]=Title, [1..12]=content, [13]=Template(3.12), [14]=Template(3.5), [15]=End

# Now move End slide to just after last content slide (position 13)
lst = prs.slides._sldIdLst
# End is at the very end (last element), move it before the two template slides
end_el = lst[-1]  # End
lst.remove(end_el)
lst.insert(len(lst) - 2, end_el)  # Insert before the two templates
# Now: [0]=Title, [1..12]=content, [13]=End, [14]=Template(3.12), [15]=Template(3.5)

# ── Save ──────────────────────────────────────────────────────────────
prs.save(DST)
print(f"\nSaved: {DST}")

# ── Verify ────────────────────────────────────────────────────────────
prs2 = Presentation(DST)
print(f"\nVerification ({len(prs2.slides)} slides):")
for i, slide in enumerate(prs2.slides):
    t = ""
    for shape in slide.shapes:
        if 'Title' in shape.name and shape.has_text_frame:
            t = shape.text_frame.text.strip()
            if t:
                break
    layout_short = slide.slide_layout.name.split(' - ')[0] if ' - ' in slide.slide_layout.name else slide.slide_layout.name
    print(f"  {i+1:2d}. [{layout_short:5s}] {t}")

print("\nNote: Slides 15-16 are leftover blank template slides — delete them in PowerPoint.")
