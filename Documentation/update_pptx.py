"""
Update LundaLogger_2026-02-27.pptx with content from PresentationOutline.md
"""
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import copy

PPTX_PATH = r'c:\_arb\VSC_projects\LundaLoggern\Documentation\LundaLogger_2026-02-27.pptx'
OUTPUT_PATH = PPTX_PATH  # overwrite in place

prs = Presentation(PPTX_PATH)

# ── Helper functions ──────────────────────────────────────────────────

def clear_placeholder(shape):
    """Remove all paragraphs from a text frame, leaving one empty paragraph."""
    tf = shape.text_frame
    for para in tf.paragraphs:
        for run in para.runs:
            run.text = ""
    # Clear using XML approach: remove all <a:p> then add one empty
    from lxml import etree
    txBody = tf._txBody
    nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
    for p in txBody.findall('a:p', nsmap):
        txBody.remove(p)
    # Add one empty paragraph back
    new_p = etree.SubElement(txBody, '{http://schemas.openxmlformats.org/drawingml/2006/main}p')


def add_bullet(text_frame, text, level=0, bold=False, font_size=None, is_first=False):
    """Add a bullet paragraph to a text frame."""
    if is_first and not text_frame.paragraphs[0].text.strip():
        para = text_frame.paragraphs[0]
    else:
        para = text_frame.add_paragraph()
    
    para.level = level
    run = para.add_run()
    run.text = text
    if bold:
        run.font.bold = True
    if font_size:
        run.font.size = Pt(font_size)


def set_content(slide, placeholder_name, bullets, font_size=None):
    """Find a placeholder by name and fill it with bullet content."""
    shape = None
    for s in slide.shapes:
        if s.name == placeholder_name and s.has_text_frame:
            shape = s
            break
    if shape is None:
        print(f"  WARNING: Could not find placeholder '{placeholder_name}'")
        return
    
    clear_placeholder(shape)
    tf = shape.text_frame
    tf.word_wrap = True
    
    first = True
    for item in bullets:
        if isinstance(item, str):
            add_bullet(tf, item, level=0, font_size=font_size, is_first=first)
            first = False
        elif isinstance(item, tuple):
            text, level, bold = item
            add_bullet(tf, text, level=level, bold=bold, font_size=font_size, is_first=first)
            first = False


def get_slide(idx):
    """Get slide by 1-based index."""
    return prs.slides[idx - 1]


def delete_slides(slide_indices):
    """Delete slides by 1-based indices (descending order to preserve indices)."""
    from lxml import etree
    for idx in sorted(slide_indices, reverse=True):
        rId = prs.slides._sldIdLst[idx - 1].get(
            '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id'
        )
        prs.part.drop_rel(rId)
        slide_elem = prs.slides._sldIdLst[idx - 1]
        prs.slides._sldIdLst.remove(slide_elem)


# ── Slide 1: Title Slide ─────────────────────────────────────────────
print("Updating Slide 1: Title")
slide1 = get_slide(1)
for shape in slide1.shapes:
    if shape.name == "Title 1" and shape.has_text_frame:
        clear_placeholder(shape)
        tf = shape.text_frame
        run = tf.paragraphs[0].add_run()
        run.text = "LundaLoggern"
    elif shape.name == "Subtitle 2" and shape.has_text_frame:
        clear_placeholder(shape)
        tf = shape.text_frame
        run = tf.paragraphs[0].add_run()
        run.text = "Smart Data Logging for Intensive Care Ventilators"
    elif shape.name == "Text Placeholder 3" and shape.has_text_frame:
        clear_placeholder(shape)
        tf = shape.text_frame
        run = tf.paragraphs[0].add_run()
        run.text = "Åke — 2026-03-16"

# ── Slide 4: The Problem ─────────────────────────────────────────────
print("Updating Slide 4: The Problem")
slide4 = get_slide(4)
set_content(slide4, "Content Placeholder 2", [
    "Intensive care ventilators generate critical patient data — but capturing it consistently is hard.",
    "Manual data recording is error-prone and time-consuming for clinical staff.",
    "Existing solutions are often expensive, complex to set up, or tied to proprietary hospital IT systems.",
    ("Key question: How can we make ventilator data logging simple, portable, and independent?", 0, True),
])

# ── Slide 5: Introducing LundaLoggern ────────────────────────────────
print("Updating Slide 5: Introducing LundaLoggern")
slide5 = get_slide(5)
set_content(slide5, "Content Placeholder 2", [
    "A small, self-contained plug-and-play data logger for SERVO ventilators.",
    "Connects directly to the ventilator — no hospital network or IT integration needed.",
    "Automatically records breath-by-breath data, ventilator settings, and waveforms to an SD card.",
    "Built-in display shows live status; built-in WiFi lets you download data from your phone or laptop.",
])

# ── Slide 6: How It Works ────────────────────────────────────────────
print("Updating Slide 6: How It Works")
slide6 = get_slide(6)
set_content(slide6, "Content Placeholder 2", [
    ("Step 1: Plug LundaLoggern into the ventilator's CIE port.", 0, True),
    ("Step 2: The device auto-detects the ventilator and starts logging.", 0, True),
    ("Step 3: Press a button to enable WiFi — scan the QR code on the display.", 0, True),
    ("Step 4: Open the web dashboard on your phone/laptop to download or review data files.", 0, True),
])

# ── Slide 7: What Gets Logged ────────────────────────────────────────
print("Updating Slide 7: What Gets Logged")
slide7 = get_slide(7)
set_content(slide7, "Content Placeholder 2", [
    ("Breath metrics", 0, True),
    ("Tidal volume, respiratory rate, minute ventilation, FiO\u2082, PEEP, peak pressure, and more.", 1, False),
    ("Ventilator settings", 0, True),
    ("Mode, patient category, compliance compensation, I:E ratio.", 1, False),
    ("Waveforms", 0, True),
    ("Flow, airway pressure, volume, CO\u2082, Edi — fully configurable.", 1, False),
    "All data is timestamped automatically using the ventilator's own clock.",
    "Configurable — choose exactly which parameters to track via simple text files.",
])

# ── Slide 8: Key Features & Benefits ─────────────────────────────────
print("Updating Slide 8: Key Features & Benefits")
slide8 = get_slide(8)
features = [
    ("Plug-and-play — No setup, training, or IT support needed", 0, False),
    ("WiFi off by default — Minimises RF interference; manually enabled only when needed", 0, False),
    ("Portable & compact — Fits in a pocket; move between ventilators", 0, False),
    ("SD card storage — Standard, removable, unlimited capacity", 0, False),
    ("WiFi & web dashboard — Download data wirelessly from any device", 0, False),
    ("QR code connectivity — Connect to WiFi in seconds", 0, False),
    ("Configurable logging — Adapt to research or clinical needs without software changes", 0, False),
    ("Live display — See connection status, ventilator ID, and timestamps at a glance", 0, False),
    ("No cloud dependency — Data stays local; important for patient privacy", 0, False),
]
set_content(slide8, "Content Placeholder 2", features, font_size=12)

# ── Slide 9: Target Markets & Use Cases ──────────────────────────────
print("Updating Slide 9: Target Markets & Use Cases")
slide9 = get_slide(9)
set_content(slide9, "Content Placeholder 2", [
    ("Clinical / Neonatal ICU", 0, True),
    ("Continuous monitoring and documentation for quality assurance.", 1, False),
    ("Research data collection during clinical studies.", 1, False),
    ("Biomedical Engineering", 0, True),
    ("Ventilator performance verification and troubleshooting.", 1, False),
    ("Long-term trend analysis during service calls.", 1, False),
    ("Medical Device R&D", 0, True),
    ("Bench testing and field validation of ventilator prototypes.", 1, False),
    ("Rapid data capture without proprietary software.", 1, False),
    ("Education & Training", 0, True),
    ("Teaching tool for respiratory therapy programs.", 1, False),
    ("Live data demonstration during ventilator training sessions.", 1, False),
])

# ── Slide 10: Competitive Advantages ─────────────────────────────────
print("Updating Slide 10: Competitive Advantages")
slide10 = get_slide(10)
set_content(slide10, "Content Placeholder 2", [
    "No recurring costs — no licenses, subscriptions, or cloud fees.",
    "Vendor-independent data — plain text files, open and portable.",
    "No hospital IT involvement — standalone device with its own WiFi.",
    "Privacy by design — data never leaves the device unless the user downloads it.",
    "Low cost hardware — built on widely available ESP32-S3 platform.",
    "Fast deployment — working in under a minute, zero configuration required.",
])

# ── Slide 11: Product Demo / Live Screenshots ────────────────────────
print("Updating Slide 11: Product Demo / Live Screenshots")
slide11 = get_slide(11)
set_content(slide11, "Content Placeholder 2", [
    "Photo or live demo of the device connected to a ventilator.",
    ("TFT Display shows:", 0, True),
    ("Ventilator type + serial number", 1, False),
    ("SD and COM status indicators", 1, False),
    ("WiFi IP address", 1, False),
    ("Web dashboard on phone:", 0, True),
    ("File listing", 1, False),
    ("Download and delete buttons", 1, False),
    ("Configuration viewer", 1, False),
])

# ── Slide 12: Technical Snapshot ──────────────────────────────────────
print("Updating Slide 12: Technical Snapshot")
slide12 = get_slide(12)
set_content(slide12, "Content Placeholder 2", [
    ("Hardware: LilyGo T-Display S3 (ESP32-S3), SD card reader, RS232 interface", 0, True),
    ("Display: 320 × 170 px color TFT", 0, True),
    ("Connectivity: WiFi access point (no internet needed)", 0, True),
    ("Protocol: SERVO CIE interface (standard on SERVO ventilators)", 0, True),
    ("Power: USB-C (can be powered from the ventilator's USB port)", 0, True),
    ("Software: Open-source Arduino/PlatformIO firmware", 0, True),
])

# ── Slide 13: Roadmap / Future Possibilities ─────────────────────────
print("Updating Slide 13: Roadmap / Future Possibilities")
slide13 = get_slide(13)
set_content(slide13, "Content Placeholder 2", [
    "WiFi auto-shutoff — automatically turn off WiFi if no one connects within 5 min.",
    "Bluetooth connectivity for direct mobile app integration.",
    "Multi-ventilator support — expand beyond SERVO to other brands.",
    "Cloud sync option (opt-in) for centralized research data collection.",
    "Real-time trend graphs on the web dashboard.",
    "Battery operation for transport and ambulance use.",
    "Regulatory pathway — CE marking / FDA clearance for clinical deployment.",
])

# ── Slide 14: Call to Action ─────────────────────────────────────────
print("Updating Slide 14: Call to Action")
slide14 = get_slide(14)
set_content(slide14, "Content Placeholder 2", [
    ('"We have a working prototype — let\'s bring it to market."', 0, True),
    "",
    ("What we need from marketing:", 0, True),
    ("Market sizing and customer segmentation feedback", 1, False),
    ("Input on branding, naming, and visual identity", 1, False),
    ("Channel strategy — direct sales, distributors, OEM partnerships?", 1, False),
    ("Pricing model validation", 1, False),
])

# ── Slide 15: Q&A ───────────────────────────────────────────────────
print("Updating Slide 15: Q&A")
slide15 = get_slide(15)
# Fix the title (had a trailing newline)
for shape in slide15.shapes:
    if shape.name == "Title 3" and shape.has_text_frame:
        clear_placeholder(shape)
        tf = shape.text_frame
        run = tf.paragraphs[0].add_run()
        run.text = "Q&A"

set_content(slide15, "Content Placeholder 2", [
    "Open floor for questions.",
    "",
    ("Backup slides available with:", 0, True),
    ("Detailed configuration file examples", 1, False),
    ("Competitive landscape comparison", 1, False),
    ("Bill of materials / cost breakdown", 1, False),
])

# ── Delete empty slides 16-20 ────────────────────────────────────────
print("Deleting empty slides 16-20...")
delete_slides([16, 17, 18, 19, 20])

# ── Save ──────────────────────────────────────────────────────────────
prs.save(OUTPUT_PATH)
print(f"\nDone! Saved to: {OUTPUT_PATH}")
print(f"Total slides now: {len(prs.slides)}")
